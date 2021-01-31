import streamlit as st
import spacy
from spacy import displacy
import pandas as pd
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
from io import StringIO

import os 

SPACY_MODEL_NAMES = ["en_core_web_sm", "en_core_web_md", "de_core_news_sm"]
DEFAULT_TEXT = "Mark Elliot Zuckerberg is an American media magnate, internet entrepreneur, and philanthropist. He is known for co-founding Facebook, Inc. and serves as its chairman, chief executive officer, and controlling shareholder. He also is a co-founder of the solar sail spacecraft development project Breakthrough Starshot and serves as one of its board members. Born in White Plains, New York, Zuckerberg attended Harvard University, where he launched the Facebook social networking service from his dormitory room on February 4, 2004, with college roommates Eduardo Saverin, Andrew McCollum, Dustin Moskovitz, and Chris Hughes. Originally launched to select college campuses, the site expanded rapidly and eventually beyond colleges, reaching one billion users by 2012. Zuckerberg took the company public in May 2012 with majority shares. His net worth is estimated to be nearly $54 billion as of March 2020. In 2007, at age 23, he became the worlds youngest self-made billionaire. As of 2019, he is the only person under 50 in the Forbes ten richest people list, and the only one under 40 in the Top 20 Billionaires list"
HTML_WRAPPER = """<div style="overflow-x: auto; border: 1px solid #e6e9ef; border-radius: 0.25rem; padding: 1rem; margin-bottom: 2.5rem">{}</div>"""


# @st.cache(ignore_hash=True)
def load_model(name):
    return spacy.load(name)


# @st.cache(ignore_hash=True)
def process_text(model_name, text):
    nlp = load_model(model_name)
    return nlp(text)

def file_selector(folder_path='./data'):
    filenames = os.listdir(folder_path)
    selected_filename = st.selectbox('Select a file', filenames)
    return os.path.join(folder_path, selected_filename)

def convert_pdf_to_txt(file, isPage):
    ans = []
    file = open(file,"rb")
    rsrcmgr = PDFResourceManager()
    retstr = StringIO()
    codec = 'utf-8'
    laparams = LAParams()
    device = TextConverter(rsrcmgr, retstr, laparams=laparams)

    interpreter = PDFPageInterpreter(rsrcmgr, device)
    password = ""
    caching = True
    pagenos = set()
    count = 1
    text = None
    name = file.name.split(".")[0]
    for page in PDFPage.get_pages(file, pagenos, password=password,caching=caching, check_extractable=True):
        interpreter.process_page(page)

        if isPage:
            text = retstr.getvalue()
            ans.append(text)
            count+=1
            retstr.truncate(0)
            retstr.seek(0)
    if not isPage:
        ans = retstr.getvalue()
        # open(f"{name}.txt","w").write(text)
    device.close()
    retstr.close()
    return ans

st.sidebar.title("Interactive spaCy visualizer")
st.sidebar.markdown(
    """
Process text with [spaCy](https://spacy.io) models and visualize named entities,
dependencies and more. Uses spaCy's built-in
[displaCy](http://spacy.io/usage/visualizers) visualizer under the hood.
"""
)

user_input = st.text_input("label goes here", './data/')

st.header("EXTRACTED INFORMATION")
if(user_input[-4:]=='.pdf') :
    c = convert_pdf_to_txt(user_input, False)
    st.write(c)
    DEFAULT_TEXT=c

spacy_model = st.sidebar.selectbox("Model name", SPACY_MODEL_NAMES)
model_load_state = st.info(f"Loading model '{spacy_model}'...")
nlp = load_model(spacy_model)
model_load_state.empty()

text = st.text_area("Text to analyze", DEFAULT_TEXT)
doc = process_text(spacy_model, text)

# if "parser" in nlp.pipe_names:
#     st.header("Dependency Parse & Part-of-speech tags")
#     st.sidebar.header("Dependency Parse")
#     split_sents = st.sidebar.checkbox("Split sentences", value=True)
#     collapse_punct = st.sidebar.checkbox("Collapse punctuation", value=True)
#     collapse_phrases = st.sidebar.checkbox("Collapse phrases")
#     compact = st.sidebar.checkbox("Compact mode")
#     options = {
#         "collapse_punct": collapse_punct,
#         "collapse_phrases": collapse_phrases,
#         "compact": compact,
#     }
#     docs = [span.as_doc() for span in doc.sents] if split_sents else [doc]
#     for sent in docs:
#         html = displacy.render(sent, options=options)
#         # Double newlines seem to mess with the rendering
#         html = html.replace("\n\n", "\n")
#         if split_sents and len(docs) > 1:
#             st.markdown(f"> {sent.text}")
#         st.write(HTML_WRAPPER.format(html), unsafe_allow_html=True)

if "ner" in nlp.pipe_names:
    st.header("IMPORTANT WORDS")
    st.sidebar.header("Named Entities")
    default_labels = ["PERSON", "ORG", "GPE", "LOC", "QUANTITY", "PERCENT", "MONEY", "LAW", "FAC", "LANGUAGE", "EVENT", "DATE"]
    labels = st.sidebar.multiselect(
        "Entity labels", nlp.get_pipe("ner").labels, default_labels
    )
    html = displacy.render(doc, style="ent", options={"ents": labels})
    # Newlines seem to mess with the rendering
    html = html.replace("\n", " ")
    st.write(HTML_WRAPPER.format(html), unsafe_allow_html=True)
    attrs = ["text", "label_", "start", "end", "start_char", "end_char"]
    if "entity_linker" in nlp.pipe_names:
        attrs.append("kb_id_")
    data = [
        [str(getattr(ent, attr)) for attr in attrs]
        for ent in doc.ents
        if ent.label_ in labels
    ]
    df = pd.DataFrame(data, columns=attrs)
    st.dataframe(df)


if "textcat" in nlp.pipe_names:
    st.header("Text Classification")
    st.markdown(f"> {text}")
    df = pd.DataFrame(doc.cats.items(), columns=("Label", "Score"))
    st.dataframe(df)

from spacy.lang.en.stop_words import STOP_WORDS
from heapq import nlargest

def text_summarizer(raw_docx, level):
    raw_text = raw_docx
    docx = nlp(raw_text)
    stopwords = list(STOP_WORDS)
    # Build Word Frequency
# word.text is tokenization in spacy
    word_frequencies = {}  
    for word in docx:  
        if word.text not in stopwords:
            if word.text not in word_frequencies.keys():
                word_frequencies[word.text] = 1
            else:
                word_frequencies[word.text] += 1


    maximum_frequncy = max(word_frequencies.values())

    for word in word_frequencies.keys():  
        word_frequencies[word] = (word_frequencies[word]/maximum_frequncy)
    # Sentence Tokens
    sentence_list = [ sentence for sentence in docx.sents ]

    # Calculate Sentence Score and Ranking
    sentence_scores = {}  
    for sent in sentence_list:  
        for word in sent:
            if word.text.lower() in word_frequencies.keys():
                if len(sent.text.split(' ')) < 30:
                    if sent not in sentence_scores.keys():
                        sentence_scores[sent] = word_frequencies[word.text.lower()]
                    else:
                        sentence_scores[sent] += word_frequencies[word.text.lower()]

    # Find N Largest
    summary_sentences = nlargest(level, sentence_scores, key=sentence_scores.get)
    final_sentences = [ w.text for w in summary_sentences ]
    summary = ' '.join(final_sentences)
    return summary

st.header("SUMMARY")

summary_level = st.text_input("Threshold Summary", 10)
st.write(text_summarizer(DEFAULT_TEXT, int(summary_level)))

st.header("Automatic PPT Generation")

ppt_data = {
        "title":"What is Global Warming",
        "slides":[{
                "title":"Introduction",
                "data":[
                    "Glaciers are melting, sea levels are rising, cloud forests are dying, and wildlife is scrambling to keep pace.",
                    "Humans have caused most of the past century's warming by releasing heat-trapping gases as we power our modern lives.",
                    "This is causing a set of changes to the Earth's climate",
                ]
            },
            {
                "title":"Understanding the greenhouse effect",
                "data":[
                    "The \"greenhouse effect\" is the warming that happens when certain gases in Earth's atmosphere trap heat.",
                    "Sunlight shines onto the Earth's surface, where the energy is absorbed and then radiate back into the atmosphere as heat.",
                    "Heat gets locked up in these molecules",
                ]
            },
            {
                "title":"Aren't temperature changes natural?",
                "data":[
                    "Human activity isn't the only factor that affects Earth's climate.",
                    "Volcanic eruptions and variations in solar radiation from sunspots etc also play a role.",
                    "But climate models that scientists use to monitor Earth’s temperatures take those factors into account",
                ]
            }]
    }

from pptx import Presentation
def ppt_generator(data):
    title_text = data.get('title','Title')
    slides = data.get('slides',[])
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    for slide_data in slides:
        bullet_slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = slide_data['title']

        tf = body_shape.text_frame
        data = slide_data['data']
        tf.text = data[0]

        for text in data[1:]:
            p = tf.add_paragraph()
            p.text = text
            p.level = 0

    prs.save('test.pptx')
    href = f'<a href="http://localhost:8000/test.pptx">Download file</a>'
    st.markdown(href, unsafe_allow_html=True)

if st.button("Generate"):
    ppt_generator(ppt_data)


st.header("QUESTION/ANSWERING")

user_input2 = st.text_input("Ask your Questions..", 'When does global warming occur?')

from allennlp.predictors.predictor import Predictor
predictor = Predictor.from_path("https://storage.googleapis.com/allennlp-public-models/bidaf-elmo-model-2018.11.30-charpad.tar.gz")
prediction = predictor.predict(
passage=DEFAULT_TEXT,
question=user_input2
)

st.write(prediction['best_span_str'])



vector_size = nlp.meta.get("vectors", {}).get("width", 0)
if vector_size:
    st.header("Vectors & Similarity")
    st.code(nlp.meta["vectors"])
    text1 = st.text_input("Text or word 1", "apple")
    text2 = st.text_input("Text or word 2", "orange")
    doc1 = process_text(spacy_model, text1)
    doc2 = process_text(spacy_model, text2)
    similarity = doc1.similarity(doc2)
    if similarity > 0.5:
        st.success(similarity)
    else:
        st.error(similarity)

st.header("Token attributes")

if st.button("Show token attributes"):
    attrs = [
        "idx",
        "text",
        "lemma_",
        "pos_",
        "tag_",
        "dep_",
        "head",
        "ent_type_",
        "ent_iob_",
        "shape_",
        "is_alpha",
        "is_ascii",
        "is_digit",
        "is_punct",
        "like_num",
    ]
    data = [[str(getattr(token, attr)) for attr in attrs] for token in doc]
    df = pd.DataFrame(data, columns=attrs)
    st.dataframe(df)


st.header("JSON Doc")
if st.button("Show JSON Doc"):
    st.json(doc.to_json())

st.header("JSON model meta")
if st.button("Show JSON model meta"):
    st.json(nlp.meta)

